#!/usr/bin/env python
"""
extract_pptx.py — PPT → data.js 提取脚本（支持增量更新）
============================================================

对应 SKILL.md 的 Skill 1（文本解析）+ Skill 2（表格提取）+ Skill 3（数据平铺）完整实现。

功能:
  1. 扫描 ppts/ 目录下的所有 .pptx 文件
  2. 解析每个 PPT 中的 13 列表格（跳过前 2 行表头）
  3. 对列 7-11 的投资者文本执行正则解析
  4. 平铺为"一笔投资一条记录"的扁平数组
  5. 按定价日期降序排列后输出为 data/data.js

增量模式（默认）:
  - 通过 data/.pptx_index.json 追踪已处理的 PPT 文件
  - 每次只提取新增或修改过的 PPT
  - 新数据与历史数据合并去重
  - 同时输出 data/data.js（前端用）和 data/data.json（Python 回读用）

用法:
  python scripts/extract_pptx.py                              # 默认模式（自动检测增量）
  python scripts/extract_pptx.py --force                      # 强制全量重新提取
  python scripts/extract_pptx.py --incremental                # 显式指定增量模式
  python scripts/extract_pptx.py --input-dir ppts             # 指定输入目录
  python scripts/extract_pptx.py --input 2026基石_0622.pptx   # 指定单个文件
  python scripts/extract_pptx.py --output data/data.js        # 指定输出路径

依赖:
  pip install python-pptx

规范依据:
  - references.md §1: PPT 表格结构映射
  - references.md §2: 文本解析边界条件
  - references.md §3: 格式规范（日期/股票代码/金额）
  - scripts.md §2: extract_pptx.py 详细规范
"""

import re
import os
import sys
import json
import argparse
import hashlib
from datetime import datetime
from pptx import Presentation


# =============================================================================
# 常量定义
# =============================================================================

# 列 7-11 对应的投资者类别中文标签
INVESTOR_CATEGORIES = {
    7: "上下游企业",
    8: "政府基金",
    9: "老股东",
    10: "战略合作",
    11: "财务投资者",
}

# 投资者名称解析正则
# 说明:
#   - (.+)        贪婪匹配投资者名称（尽可能多，确保匹配到最后一对括号作为金额）
#   - [（(]       匹配全角或半角左括号（金额的开始）
#   - \s*         括号内可能有的空格
#   - ([\d,]+\.?\d*)  金额数字，含千分位逗号和小数
#   - \s*         括号内可能有的空格
#   - [)）]       匹配全角或半角右括号
# 用例覆盖（references.md §2）:
#   A: Oakwise (9.57)           → name="Oakwise",         amount=9.57
#   B: 上海闵行 (3)              → name="上海闵行",         amount=3.0
#   C: Tembusu (David Su) (2.04) → name="Tembusu (David Su)", amount=2.04
#   D: 奥博资本 Orbimed (18)     → name="奥博资本 Orbimed",  amount=18.0
#   F: 欣旺达香港（5.20）         → name="欣旺达香港",        amount=5.20
INVESTOR_RE = re.compile(r'^(.+)[（(]\s*([\d,]+\.?\d*)\s*[)）]$')

# 日期格式统一（兼容两种写法）
DATE_RE = re.compile(r'^(\d{4})[-/](\d{1,2})[-/](\d{1,2})$')

# 需要跳过单元格内容
SKIP_VALUES = {"", "-", "—", "/", "na", "N/A", "n/a"}

# 投资者别名规范化配置
ALIASES_PATH = os.path.join(os.path.dirname(os.path.dirname(__file__)), "data", "investor_aliases.json")

# 增量更新相关路径
PPTX_INDEX_PATH = os.path.join(os.path.dirname(os.path.dirname(__file__)), "data", ".pptx_index.json")
DATA_JSON_PATH = os.path.join(os.path.dirname(os.path.dirname(__file__)), "data", "data.json")


def load_aliases() -> dict:
    """加载投资者别名规范化规则。"""
    if os.path.exists(ALIASES_PATH):
        with open(ALIASES_PATH, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}


def normalize_investor_name(name: str, aliases: dict) -> str:
    """
    对投资者名称做规范化清洗。

    规则来源: data/investor_aliases.json
      - exact: 精确匹配替换
      - prefix_remove: 删除指定前缀
      - case_insensitive: 不区分大小写匹配后替换
    """
    if not name:
        return name

    # 1. 精确匹配替换
    for old, new in aliases.get("exact", {}).items():
        if name == old:
            return new

    # 1b. 包含关键字匹配（不区分大小写）
    for keyword, new in aliases.get("contains", {}).items():
        if keyword.lower() in name.lower():
            return new

    # 2. 子串替换（如删除名称中间的"通过"）
    for rule in aliases.get("substring", []):
        name = name.replace(rule["old"], rule["new"])
    name = name.strip()

    # 3. 不区分大小写匹配替换
    name_lower = name.lower()
    for old_lower, new in aliases.get("case_insensitive", {}).items():
        if name_lower == old_lower.lower():
            return new

    return name


# =============================================================================
# 核心函数
# =============================================================================


def normalize_date(raw: str) -> str:
    """
    统一日期格式为 YYYY-MM-DD。

    兼容:
      - 2026-06-17  → 2026-06-17
      - 2026/5/6    → 2026-05-06
    """
    match = DATE_RE.match(raw.strip())
    if not match:
        return raw.strip()  # 无法解析则原样返回
    year, month, day = match.groups()
    return f"{year}-{int(month):02d}-{int(day):02d}"


def normalize_stock_code(raw: str) -> str:
    """
    统一股票代码格式，补上 .HK 后缀。

    兼容:
      - 6675     → 6675.HK
      - 9981.HK  → 9981.HK
      - 0068     → 0068.HK
    """
    code = raw.strip()
    if not code:
        return code
    if code.endswith(".HK"):
        return code
    # 纯数字或有前缀字母的代码，补 .HK
    return f"{code}.HK"


def parse_number(raw: str) -> float | None:
    """
    解析数字字符串，去千分位逗号，转浮点数。

    示例:
      "981"    → 981.0
      "2,567"  → 2567.0
      "5.20"   → 5.20
      ""       → None
    """
    s = raw.strip().replace(",", "").replace("，", "")
    if not s:
        return None
    try:
        return float(s)
    except ValueError:
        return None


def parse_investors(cell_text: str, category: str) -> list[dict]:
    """
    解析单元格内的投资者文本（Skill 1）。

    输入: "Oakwise (9.57)\\nTembusu (David Su) (2.04)\\n阎焱 (2.04)"
    输出: [
            {"name": "Oakwise", "amount_mn": 9.57, "category": "财务投资者"},
            {"name": "Tembusu (David Su)", "amount_mn": 2.04, "category": "财务投资者"},
            ...
          ]

    处理:
      - 以 \n 和 \x0b 为候选分隔符切分行（\x0b → \n 统一处理）
      - 每行尝试匹配正则（支持全角/半角括号）
      - 不匹配的行缓存起来，向前合并到下一匹配行一起解析
        （解决 PPT 手工换行导致投资者名被截断的问题）
    """
    if not cell_text:
        return []

    # 归一化：\x0b（PPT 垂直制表符）→ \n，\xa0（不换行空格）→ 普通空格
    text = cell_text.replace("\x0b", "\n").replace("\xa0", " ")

    results = []
    buffer = ""  # 缓存不匹配的行，向前合并
    for line in text.split("\n"):
        line = line.strip()
        if not line or line in SKIP_VALUES:
            continue

        match = INVESTOR_RE.match(line)
        if match:
            # 此行匹配 → 如果有缓存，先合并到名称前
            name = ((buffer + match.group(1)) if buffer else match.group(1)).strip()
            amount_str = match.group(2).replace(",", "")
            buffer = ""  # 清空缓存
            try:
                amount = float(amount_str)
                results.append({
                    "investor_name": name,
                    "investor_category": category,
                    "amount_mn": amount,
                })
            except ValueError:
                pass
        else:
            # 此行不匹配 → 缓存起来（可能是上一行的续行）
            buffer += line + " "

    return results


def extract_pptx(filepath: str) -> list[dict]:
    """
    从单个 PPTX 文件中提取所有投资记录（Skill 2 + Skill 3）。

    返回: 按 (pricing_date DESC, company_name ASC) 排序的记录列表
    """
    filename = os.path.basename(filepath)
    prs = Presentation(filepath)

    all_records = []
    aliases = load_aliases()

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_table:
                continue

            table = shape.table
            n_rows = len(table.rows)
            n_cols = len(table.columns)

            # 安全校验：至少 3 行（2 行表头 + 1 行数据）且 12+ 列
            if n_rows < 3 or n_cols < 12:
                continue

            # 从第 3 行开始遍历数据行（索引 2，跳过 2 行表头）
            for r_idx in range(2, n_rows):
                row = table.rows[r_idx]

                # ---------- 读取项目基本信息（列 0-6） ----------
                pricing_date = normalize_date(row.cells[0].text)
                stock_code = normalize_stock_code(row.cells[1].text)
                company_name = row.cells[2].text.strip()
                # 行业：PPT 中人工换行不做子类分隔，直接合并为连续字符串（无空格）
                raw_industry = row.cells[3].text.strip()
                industry = "".join(raw_industry.replace("\x0b", "").replace("\n", "").replace("\r", "").split())
                if not industry:
                    industry = None

                ipo_size = parse_number(row.cells[4].text)
                cs_size = parse_number(row.cells[5].text)
                cs_ratio = row.cells[6].text.strip() or None

                # 跳过空行（无公司名视为空行）
                if not company_name and not stock_code:
                    continue

                # ---------- 读取投资者信息（列 7-11，平铺输出） ----------
                for col in range(7, 12):
                    cell_text = row.cells[col].text
                    category = INVESTOR_CATEGORIES.get(col, "未知")
                    investors = parse_investors(cell_text, category)

                    for inv in investors:
                        # 投资者名称规范化清洗
                        inv["investor_name"] = normalize_investor_name(inv["investor_name"], aliases)
                        record = {
                            # 项目信息
                            "pricing_date": pricing_date,
                            "stock_code": stock_code,
                            "company_name": company_name,
                            "industry": industry,
                            "ipo_size_mn": ipo_size,
                            "cs_size_mn": cs_size,
                            "cs_ratio": cs_ratio,
                            # 该笔投资信息
                            "investor_name": inv["investor_name"],
                            "investor_category": inv["investor_category"],
                            "amount_mn": inv["amount_mn"],
                            # 来源追溯
                            "source_ppt": filename,
                            "slide_no": slide_idx,
                        }
                        all_records.append(record)

    return all_records


# =============================================================================
# 增量更新相关函数
# =============================================================================


def make_record_key(record: dict) -> tuple:
    """
    生成记录的自然键，用于去重。

    同一投资者对同一项目的投资金额应为唯一。
    """
    return (
        record.get("stock_code", "") or "",
        record.get("investor_name", "") or "",
        record.get("amount_mn", 0) or 0,
    )


def get_file_signature(filepath: str) -> dict:
    """
    获取文件签名（SHA-256 + 大小），用于检测文件是否变化。

    使用内容哈希而不是 mtime，确保 GitHub Actions / Netlify 等不同
    运行环境中不会因为文件时间戳变化而把所有历史 PPT 误判为已修改。
    """
    stat = os.stat(filepath)
    digest = hashlib.sha256()
    with open(filepath, "rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            digest.update(chunk)
    return {
        "sha256": digest.hexdigest(),
        "size": stat.st_size,
    }


def load_pptx_index() -> dict:
    """
    加载已处理的 PPT 文件索引。

    结构:
    {
        "version": 1,
        "processed_files": {
            "2026基石_0622.pptx": { "mtime": 1234567890.0, "size": 1315835, "record_count": 494 }
        },
        "last_updated": "2026-06-24 22:30:36"
    }
    """
    if os.path.exists(PPTX_INDEX_PATH):
        try:
            with open(PPTX_INDEX_PATH, "r", encoding="utf-8") as f:
                return json.load(f)
        except (json.JSONDecodeError, Exception):
            pass
    return {"version": 1, "processed_files": {}, "last_updated": None}


def save_pptx_index(pptx_files: list[str], record_counts: dict[str, int] | None = None):
    """
    保存已处理的 PPT 文件索引。

    Args:
        pptx_files: 已成功提取的 PPT 文件路径列表
        record_counts: 可选，每个文件的记录数映射
    """
    index = load_pptx_index()
    processed = index.setdefault("processed_files", {})

    for f in pptx_files:
        basename = os.path.basename(f)
        sig = get_file_signature(f)
        processed[basename] = {
            "size": sig["size"],
            "record_count": record_counts.get(basename, 0) if record_counts else 0,
        }

    index["last_updated"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    index["version"] = 2

    os.makedirs(os.path.dirname(PPTX_INDEX_PATH), exist_ok=True)
    with open(PPTX_INDEX_PATH, "w", encoding="utf-8") as f:
        json.dump(index, f, ensure_ascii=False, indent=2)


def find_new_or_changed_files(pptx_files: list[str], index: dict) -> list[str]:
    """
    与索引比对，找出新增或修改过的 PPT 文件。

    判定规则:
      - 文件名不在索引中 → 新文件
      - 文件名在索引中，但 mtime 或 size 不同 → 被修改过的文件
    """
    processed = index.get("processed_files", {})
    changed = []
    for f in pptx_files:
        basename = os.path.basename(f)
        sig = get_file_signature(f)
        if basename not in processed:
            changed.append(f)
        elif (processed[basename].get("sha256") != sig["sha256"] or
              processed[basename].get("size") != sig["size"]):
            changed.append(f)
    return changed


def load_existing_json_data() -> list[dict] | None:
    """
    从 data/data.json 加载已有记录（供增量合并用）。

    如果文件不存在或格式损坏，返回 None。
    """
    if os.path.exists(DATA_JSON_PATH):
        try:
            with open(DATA_JSON_PATH, "r", encoding="utf-8") as f:
                return json.load(f)
        except (json.JSONDecodeError, Exception):
            pass
    return None


def merge_records(existing: list[dict], new_records: list[dict]) -> list[dict]:
    """
    将新提取的记录与历史记录合并去重。

    去重策略:
      - 自然键: (stock_code, investor_name, amount_mn)
      - 先加载历史记录
      - 新记录如果键已存在，以新记录为准（覆盖更新）
      - 新记录如果键不存在，追加
      - 始终优先保留有新数据的 source_ppt（最新来源）
    """
    # 构建历史记录的键映射 {key: index}
    seen = {}
    merged = []

    for r in existing:
        key = make_record_key(r)
        if key not in seen:
            seen[key] = len(merged)
            merged.append(dict(r))  # 深拷贝，避免后续修改影响原始数据

    for r in new_records:
        key = make_record_key(r)
        if key in seen:
            # 用新记录覆盖（保留新数据的 source_ppt，更准确）
            merged[seen[key]] = dict(r)
        else:
            seen[key] = len(merged)
            merged.append(dict(r))

    return merged


def get_date_cutoff_from_records(records: list[dict]) -> str | None:
    """
    从记录集合中找出最晚的定价日期作为数据截止日。
    """
    valid_dates = []
    for r in records:
        d = r.get("pricing_date", "")
        if d and DATE_RE.match(d):
            valid_dates.append(d)
    if valid_dates:
        valid_dates.sort(reverse=True)
        return valid_dates[0]
    return None


def get_all_source_ppts(records: list[dict]) -> list[str]:
    """
    从记录集合中提取所有不重复的来源 PPT 文件名（按首次出现顺序）。
    """
    seen = set()
    sources = []
    for r in records:
        f = r.get("source_ppt", "")
        if f and f not in seen:
            seen.add(f)
            sources.append(f)
    return sources


# =============================================================================
# 输出格式化
# =============================================================================


def sort_records(records: list[dict], sort_order: str = "date-desc") -> list[dict]:
    """
    按定价日期降序（最新在前），同日期内按公司名升序排列。
    """
    def sort_key(r):
        d = r["pricing_date"]
        if not d or not DATE_RE.match(d or ""):
            return ("0000-00-00", r.get("company_name") or "")
        return (d, r.get("company_name") or "")

    if sort_order == "date-desc":
        return sorted(records, key=sort_key, reverse=True)
    else:
        def sort_key_asc(r):
            d = r["pricing_date"]
            if not d or not DATE_RE.match(d or ""):
                return ("9999-99-99", r.get("company_name") or "")
            return (d, r.get("company_name") or "")
        return sorted(records, key=sort_key_asc)


def records_to_js(records: list[dict], source_files: list[str], date_cutoff: str | None = None) -> str:
    """
    将记录列表格式化为 data.js 文件内容。

    输出格式（scripts.md §2.3）:
      window.cornerstoneData = [
        { ... },
        ...
      ];
    """
    lines = [
        "// data/data.js — 自动生成，请勿手动修改",
        f"// 生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
    ]
    if source_files:
        lines.append(f"// 来源文件: {', '.join(source_files)}")
    if date_cutoff:
        lines.append(f"// 数据截止: {date_cutoff}")
    lines.append("")
    lines.append("window.cornerstoneData = [")

    for i, r in enumerate(records):
        # 格式化每个记录为 JSON 行
        line = "  {"
        line += f'pricing_date: {json_dumps(r["pricing_date"])}, '
        line += f'stock_code: {json_dumps(r["stock_code"])}, '
        line += f'company_name: {json_dumps(r["company_name"])}, '
        line += f'industry: {json_dumps(r["industry"])}, '
        line += f'ipo_size_mn: {json_val(r["ipo_size_mn"])}, '
        line += f'cs_size_mn: {json_val(r["cs_size_mn"])}, '
        line += f'cs_ratio: {json_dumps(r["cs_ratio"])}, '
        line += f'investor_name: {json_dumps(r["investor_name"])}, '
        line += f'investor_category: {json_dumps(r["investor_category"])}, '
        line += f'amount_mn: {json_val(r["amount_mn"])}, '
        line += f'source_ppt: {json_dumps(r["source_ppt"])}, '
        line += f'slide_no: {r["slide_no"]}'
        line += "},"
        lines.append(line)

    lines.append("];")
    if date_cutoff:
        lines.append(f'window.dataCutoff = "{date_cutoff}";')
    lines.append("")
    return "\n".join(lines)


def json_dumps(val) -> str:
    """安全的 JSON 字符串序列化（处理 None 和特殊字符）"""
    if val is None:
        return "null"
    s = str(val)
    # 转义特殊字符
    s = s.replace("\\", "\\\\").replace('"', '\\"').replace("\n", "\\n").replace("\r", "\\r").replace("\t", "\\t")
    return f'"{s}"'


def json_val(val) -> str:
    """数值序列化"""
    if val is None:
        return "null"
    if isinstance(val, float):
        # 如果是整数浮点，输出整数形式
        if val == int(val):
            return str(int(val))
        return str(val)
    return str(val)


# =============================================================================
# 入口
# =============================================================================


def find_pptx_files(input_path: str) -> list[str]:
    """递归查找所有 .pptx 文件。"""
    if os.path.isfile(input_path):
        return [input_path]
    if os.path.isdir(input_path):
        files = []
        for f in sorted(os.listdir(input_path)):
            if f.startswith("~$"):
                continue  # 跳过 Office 临时文件
            if f.lower().endswith(".pptx"):
                files.append(os.path.join(input_path, f))
        return files
    return []


def extract_date_cutoff(files: list[str]) -> str | None:
    """
    从文件名推测数据截止日期（取最新文件的日期）。

    命名格式: 2026基石_0622.pptx → 2026-06-22
           2025基石1215.pptx → 2025-12-15
           20260622.pptx → 2026-06-22（纯数字 fallback）

    返回最新文件的日期，越新的文件优先级越高。
    """
    latest = None
    for f in files:
        basename = os.path.splitext(os.path.basename(f))[0]
        date_str = None

        # 匹配 "2026基石_0622"、"2025基石1215" 模式
        # 用 \D+ 匹配中间的汉字/分隔符，末尾 $ 确保取到最后的 MMDD
        m = re.search(r'(\d{4})\D+(\d{2})(\d{2})$', basename)
        if m:
            date_str = f"{m.group(1)}-{m.group(2)}-{m.group(3)}"
        else:
            # 纯数字 fallback: "20260622"
            m = re.search(r'(\d{4})(\d{2})(\d{2})', basename)
            if m:
                date_str = f"{m.group(1)}-{m.group(2)}-{m.group(3)}"

        if date_str and (latest is None or date_str > latest):
            latest = date_str

    return latest


def print_stats(records: list[dict], output_path: str, pptx_count: int, date_cutoff: str | None):
    """打印统计摘要。"""
    companies = set(r["company_name"] for r in records if r["company_name"])
    investors = set(r["investor_name"] for r in records if r["investor_name"])
    industries = set()
    for r in records:
        if r["industry"]:
            industries.add(r["industry"].split("/")[0].strip())
    total_amount = sum(r["amount_mn"] for r in records if r["amount_mn"])

    print()
    print(f"  DONE -> {output_path}")
    print(f"  Total records: {len(records)}")
    print(f"  Unique projects: {len(companies)}")
    print(f"  Unique investors: {len(investors)}")
    print(f"  Industries covered: {len(industries)}")
    print(f"  Total cornerstone: ${total_amount:,.2f}M")
    print(f"  Source files: {pptx_count}")
    if date_cutoff:
        print(f"  Data cutoff: {date_cutoff}")


def main():
    parser = argparse.ArgumentParser(
        description="港股 IPO 基石投资者数据提取脚本 — PPTX → data.js",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    parser.add_argument(
        "--input-dir", default="ppts",
        help="PPT 输入目录（默认: ppts/）"
    )
    parser.add_argument(
        "--input", "-i", dest="single_file",
        help="指定单个 PPT 文件（优先级高于 --input-dir）"
    )
    parser.add_argument(
        "--output", "-o", default="data/data.js",
        help="输出文件路径（默认: data/data.js）"
    )
    parser.add_argument(
        "--sort", choices=["date-desc", "date-asc"], default="date-desc",
        help="排序方式（默认: date-desc 日期降序）"
    )
    parser.add_argument(
        "--quiet", "-q", action="store_true",
        help="静默模式，不输出统计信息"
    )
    # --- 增量更新参数 ---
    mode_group = parser.add_mutually_exclusive_group()
    mode_group.add_argument(
        "--incremental", "-n", action="store_true", default=None,
        help="增量模式：只提取新增/修改的 PPT，与历史数据合并（默认自动检测）"
    )
    mode_group.add_argument(
        "--force", "-f", action="store_true",
        help="强制全量模式：重新提取所有 PPT，覆盖历史数据"
    )
    args = parser.parse_args()

    # ---- 确定输入文件 ----
    if args.single_file:
        pptx_files = find_pptx_files(args.single_file)
    else:
        pptx_files = find_pptx_files(args.input_dir)

    if not pptx_files:
        print(f"ERROR: No .pptx files found at: {args.single_file or args.input_dir}")
        sys.exit(1)

    # ---- 判断运行模式 ----
    index = load_pptx_index()
    has_index = bool(index.get("processed_files"))

    if args.force:
        mode = "full"
        if not args.quiet:
            print(f"  [MODE] 强制全量提取（--force）")
    elif args.incremental is True:
        mode = "incremental"
        if not args.quiet:
            print(f"  [MODE] 增量模式（--incremental）")
    elif args.incremental is False:
        mode = "full"
        if not args.quiet:
            print(f"  [MODE] 全量模式（--incremental=false）")
    else:
        # 自动检测：如果 data/.pptx_index.json 和 data/data.json 都存在，自动进入增量模式
        if has_index and os.path.exists(DATA_JSON_PATH):
            mode = "incremental"
            if not args.quiet:
                print(f"  [MODE] 自动增量模式（检测到已有索引和数据）")
        else:
            mode = "full"
            if not args.quiet:
                print(f"  [MODE] 全量模式（首次运行，无历史数据）")

    # ---- 增量模式：只取新文件，加载历史数据 ----
    if mode == "incremental":
        # 检查哪些文件是新/改过的
        changed_files = find_new_or_changed_files(pptx_files, index)
        if not changed_files:
            print(f"  [OK] 所有 PPT 文件已是最新，无需更新")
            return

        # 加载历史数据
        existing = load_existing_json_data()
        if existing is None:
            print(f"  [警告] 无法读取 data/data.json，自动切换为全量模式")
            mode = "full"
        else:
            if not args.quiet:
                unchanged = len(pptx_files) - len(changed_files)
                print(f"  FILES: {len(pptx_files)} total, {len(changed_files)} new/changed, {unchanged} unchanged (cached)")

            # 只提取新增/修改的文件
            all_new_records = []
            record_counts = {}
            for f in changed_files:
                records = extract_pptx(f)
                all_new_records.extend(records)
                record_counts[os.path.basename(f)] = len(records)
                if not args.quiet:
                    print(f"  [FILE] {os.path.basename(f)} -> {len(records)} records")

            if not all_new_records:
                print("  [警告] 新 PPT 中未提取到任何投资记录")
                return

            # ---- 轻量检查：新名称是否被别名覆盖（不阻塞流程） ----
            try:
                from check_new_names import check_new_names
                _aliases = load_aliases()
                check_new_names(all_new_records, _aliases, quiet=args.quiet)
            except ImportError:
                pass  # 没有 check_new_names 模块不影响主流程

            # 合并去重
            merged = merge_records(existing, all_new_records)

            if not merged:
                print("ERROR: Merge resulted in empty dataset")
                sys.exit(1)

            all_records = merged
            pptx_for_index = changed_files
            record_counts_for_index = record_counts

    # ---- 全量模式：重新提取全部 ----
    if mode == "full":
        all_records = []
        record_counts = {}
        for f in pptx_files:
            records = extract_pptx(f)
            all_records.extend(records)
            record_counts[os.path.basename(f)] = len(records)
            if not args.quiet:
                print(f"  [FILE] {os.path.basename(f)} -> {len(records)} records")

        if not all_records:
            print("ERROR: No investment records extracted")
            sys.exit(1)

        # 全量模式下，先用空索引覆盖（清除旧记录）
        with open(PPTX_INDEX_PATH, "w", encoding="utf-8") as f:
            json.dump({"version": 2, "processed_files": {}, "last_updated": None}, f, ensure_ascii=False, indent=2)
        pptx_for_index = pptx_files
        record_counts_for_index = record_counts

    # ---- 排序 ----
    all_records = sort_records(all_records, args.sort)

    # ---- 确定数据截止日期 ----
    # 优先从文件名推断，其次从数据中最晚的定价日期
    date_cutoff = extract_date_cutoff(pptx_files)
    if not date_cutoff:
        date_cutoff = get_date_cutoff_from_records(all_records)

    # ---- 收集所有来源 PPT 文件名 ----
    all_sources = get_all_source_ppts(all_records)

    # ---- 输出 data.js ----
    output_path = args.output
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

    js_content = records_to_js(
        all_records,
        source_files=all_sources,
        date_cutoff=date_cutoff,
    )

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(js_content)

    # ---- 输出 data.json（供 Python 下次增量读回） ----
    os.makedirs(os.path.dirname(DATA_JSON_PATH) or ".", exist_ok=True)
    # 将 JS 格式的记录转换为纯净 JSON（序列化为 JSON 数组）
    def serialize_record(r):
        """将记录转为纯 JSON 可序列化的 dict。"""
        return {
            "pricing_date": r["pricing_date"],
            "stock_code": r["stock_code"],
            "company_name": r["company_name"],
            "industry": r["industry"],
            "ipo_size_mn": r["ipo_size_mn"],
            "cs_size_mn": r["cs_size_mn"],
            "cs_ratio": r["cs_ratio"],
            "investor_name": r["investor_name"],
            "investor_category": r["investor_category"],
            "amount_mn": r["amount_mn"],
            "source_ppt": r["source_ppt"],
            "slide_no": r["slide_no"],
        }

    with open(DATA_JSON_PATH, "w", encoding="utf-8") as f:
        json.dump([serialize_record(r) for r in all_records], f, ensure_ascii=False, indent=2)

    # ---- 更新 PPT 索引 ----
    save_pptx_index(pptx_for_index, record_counts_for_index)

    # ---- 统计摘要 ----
    if not args.quiet:
        print_stats(all_records, output_path, len(all_sources), date_cutoff)
        if mode == "incremental":
            print(f"  Mode: incremental (+{len(pptx_for_index)} files)")


if __name__ == "__main__":
    main()
